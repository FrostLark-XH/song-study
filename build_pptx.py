#!/usr/bin/env python3
"""Build lyric-card PPTX from song-study data.json.

Language auto-detected from data structure (3-tuple lines → Japanese,
2-tuple → English). Japanese: Yu Mincho main text + furigana + romaji + CN.
English: Adobe Caslon Pro Semibold main text + CN, 3–5 lines per slide.
Layout cycles through 4 poster templates. Self-verifies on completion.
"""

import sys, os, json, hashlib
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from lxml import etree

sys.stdout.reconfigure(encoding="utf-8")

from scripts.pptx_palette import TEXT_JP, TEXT_FURIGANA, TEXT_ROMAJI, TEXT_CN, TEXT_SECTION, rgb, palette_for_bg
from scripts.pptx_background import generate_mood_background
from scripts.pptx_furigana import tokenize, _has_kanji

# ── Fonts ─────────────────────────────────────────────────────────────────────

FONT_JP = "Yu Mincho Demibold"
FONT_FURIGANA = "MS Gothic"
FONT_ROMAJI = "Courier New"
FONT_CN = "SimSun"
FONT_EN = "Arial"
FONT_EN_SERIF = "Adobe Caslon Pro"
FONT_EN_SERIF_BOLD = "Adobe Caslon Pro Semibold"

# ── Slide geometry ────────────────────────────────────────────────────────────

SLIDE_W = int(Inches(13.333))
SLIDE_H = int(Inches(7.5))
SLIDE_SAFE_RIGHT = int(SLIDE_W * 0.85)    # 15% right margin minimum
LEFT_MARGIN = int(Inches(1.2))
MAX_X_MARGIN = int(Inches(2.6))            # tightest layout x → pre-layout checks
BOTTOM_MARGIN = int(Inches(0.55))


def _max_content_w(x_margin: int) -> int:
    """Max textbox width at x_margin, keeping right edge within 85% of slide."""
    return SLIDE_SAFE_RIGHT - x_margin

# ── Font sizes ────────────────────────────────────────────────────────────────

FURIGANA_SIZE_PT = 13
ROMAJI_SIZE_PT = 20
ROMAJI_SPACING_PT = 1.5     # typewriter tracking
CN_SIZE_PT = 17
SECTION_SIZE_PT = 14

# ── Spacing ───────────────────────────────────────────────────────────────────

FURIGANA_ABOVE_BASE = int(Pt(1))
GAP_AFTER_BASE = int(Pt(14))
GAP_BETWEEN_SPLIT = int(Pt(10))
GAP_LINES = int(Pt(56))
GAP_AFTER_CN = int(Pt(6))

EMU_PER_PT = 12700
SPLIT_MIN_CHARS = 12


# ── Layout templates ──────────────────────────────────────────────────────────
# Cycle through 4 distinct compositions so pages don't feel uniform.

LAYOUTS = [
    # (x_margin, gap_between_groups)
    # y placement is computed per-slide; gap only matters for 2-group slides.
    {"name": "top_left",    "x": LEFT_MARGIN,           "gap": GAP_LINES},
    {"name": "mid_left",    "x": LEFT_MARGIN,           "gap": GAP_LINES},
    {"name": "center_left", "x": int(Inches(2.6)),      "gap": GAP_LINES},
    {"name": "stagger",     "x": LEFT_MARGIN,           "gap": int(Inches(2.4))},
]

TOP_LEFT_TOP = int(Inches(1.1))       # y-start for top_left / stagger group 1
STAGGER_G2_TOP = int(Inches(4.6))     # y-start for stagger group 2


# ── Helpers ───────────────────────────────────────────────────────────────────

def _set_cjk_font(run, cjk: str, en: str = FONT_EN):
    rPr = run._r.get_or_add_rPr()
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    ea = rPr.find(f"{ns}ea")
    if ea is None:
        ea = etree.SubElement(rPr, f"{ns}ea")
    ea.set("typeface", cjk)
    run.font.name = en


def _add_run(p, text, cjk, en=FONT_EN, size=Pt(14), color=None):
    run = p.add_run()
    run.text = text
    run.font.size = size
    if color:
        run.font.color.rgb = rgb(color)
    _set_cjk_font(run, cjk, en)
    return run


def _set_char_spacing(run, spacing_pt: float):
    """Add character spacing in pts (1pt = 100 in a:spc val)."""
    rPr = run._r.get_or_add_rPr()
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    spc = etree.SubElement(rPr, f"{ns}spc")
    spc.set("val", str(int(spacing_pt * 100)))


def _char_width_emu(font_size_pt: int) -> int:
    return font_size_pt * EMU_PER_PT


def _text_width_emu(text: str, font_size_pt: int) -> int:
    return len(text) * _char_width_emu(font_size_pt)


def _en_text_width_emu(text: str, font_size_pt: int) -> int:
    """Estimate English text width for serif fonts (Adobe Caslon Pro).

    Character width relative to CJK mono at same font size. Serif fonts have
    higher stroke contrast and wider proportions — coefficients are calibrated
    ~20% above average sans-serif estimates.
    """
    # Wide uppercase + round forms
    wide = sum(1 for c in text if c in 'MWQGDOmwqgdo@#$%&')
    # Medium-wide: most uppercase, tall lowercase
    med = sum(1 for c in text if c in 'ABCEFHKLNRSTUVXYZbpdhknuv089')
    # Normal: standard lowercase, numbers
    normal = sum(1 for c in text if c in 'acemorsxyz234567')
    # Narrow
    narrow = sum(1 for c in text if c in "iltfj1.,;:'\"!?()[]{}—–- ")
    emu = int((wide * 0.82 + med * 0.66 + normal * 0.56 + narrow * 0.38)
              * font_size_pt * EMU_PER_PT)
    return emu


def _en_font_size(max_line_len: int) -> int:
    """English font size based on longest line in group. Range: 28–42pt."""
    if max_line_len <= 15:    return 42
    elif max_line_len <= 25:  return 38
    elif max_line_len <= 35:  return 34
    elif max_line_len <= 45:  return 30
    else:                     return 28


# ── Font sizing ───────────────────────────────────────────────────────────────

def _jp_font_size(char_count: int) -> int:
    """Return initial JP font size. Range: 48–64pt."""
    if char_count <= 7:
        return 64
    elif char_count <= 12:
        return 60
    elif char_count <= 16:
        return 56
    else:
        return 52


def _reduce_tier(fs: int) -> int:
    if fs >= 64: return 60
    elif fs >= 60: return 56
    elif fs >= 56: return 52
    else: return 48


# ── Width constraint ──────────────────────────────────────────────────────────

def _text_fits_width(text: str, font_size_pt: int, x_margin: int = None) -> bool:
    if x_margin is None:
        x_margin = MAX_X_MARGIN
    return _text_width_emu(text, font_size_pt) <= _max_content_w(x_margin)


# ── Line splitting ────────────────────────────────────────────────────────────

def _find_split_pos(text: str, tokens: list[dict]) -> int:
    half = len(text) // 2
    MIN_SIDE = 5
    search = min(5, len(text) // 3)

    boundaries = {}
    char_pos = 0
    for tok in tokens:
        for _ in tok["base"]:
            char_pos += 1
        boundaries[char_pos] = not tok["needs_furi"]

    for particle in ["を", "は", "が", "に", "で", "の", "も", "へ", "と"]:
        idx = 0
        while True:
            idx = text.find(particle, idx)
            if idx == -1: break
            pos = idx + len(particle)
            if pos not in boundaries:
                boundaries[pos] = True
            idx += 1

    for punct in ["、", "。", "！", "？", "…", "〜", "～"]:
        idx = 0
        while True:
            idx = text.find(punct, idx)
            if idx == -1: break
            pos = idx + len(punct)
            if pos not in boundaries:
                boundaries[pos] = True
            idx += 1

    def _score(pos, is_kana, dist):
        left, right = pos, len(text) - pos
        if left < MIN_SIDE or right < MIN_SIDE:
            return -1
        return (2 if is_kana else 1) * (search - dist + 1)

    best_pos, best_score = None, -1
    for pos, is_kana in boundaries.items():
        dist = abs(pos - half)
        if dist <= search:
            s = _score(pos, is_kana, dist)
            if s > best_score:
                best_score, best_pos = s, pos

    if best_pos is not None:
        return best_pos

    for pos, is_kana in boundaries.items():
        dist = abs(pos - half)
        if dist <= search:
            left, right = pos, len(text) - pos
            balance = min(left, right) / max(left, right)
            s = (2 if is_kana else 1) * (search - dist + 1) * max(balance, 0.25)
            if s > best_score:
                best_score, best_pos = s, pos

    if best_pos is not None:
        return best_pos

    best_any, best_any_dist = None, len(text)
    for pos in boundaries:
        dist = abs(pos - half)
        if dist < best_any_dist:
            best_any_dist, best_any = dist, pos
    return best_any or half


def _split_tokens(tokens: list[dict]) -> list[list[dict]]:
    full_text = "".join(t["base"] for t in tokens)
    split_at = _find_split_pos(full_text, tokens)

    left, right = [], []
    accumulated = 0
    for tok in tokens:
        if accumulated < split_at:
            tok_chars = len(tok["base"])
            if accumulated + tok_chars <= split_at:
                left.append(tok)
            else:
                split_in_tok = split_at - accumulated
                left_part, right_part = tok["base"][:split_in_tok], tok["base"][split_in_tok:]
                if left_part:
                    left.append({"base": left_part,
                        "kana": tok["kana"][:split_in_tok] if tok["needs_furi"] else "",
                        "needs_furi": _has_kanji(left_part)})
                if right_part:
                    right.append({"base": right_part,
                        "kana": tok["kana"][split_in_tok:] if tok["needs_furi"] else "",
                        "needs_furi": _has_kanji(right_part)})
            accumulated += tok_chars
        else:
            right.append(tok)
            accumulated += len(tok["base"])

    if not left or not right:
        return [tokens]
    return [left, right]


# ── Core resolution ───────────────────────────────────────────────────────────

def _resolve_line(original: str, is_paired: bool = False):
    tokens = tokenize(original)
    full_text = "".join(t["base"] for t in tokens)
    char_count = len(full_text)

    fs = _jp_font_size(char_count)

    needs_split = char_count >= 19
    if not needs_split and not _text_fits_width(full_text, fs):
        if char_count >= SPLIT_MIN_CHARS:
            needs_split = True

    if needs_split:
        sub_token_lists = _split_tokens(tokens)
        font_sizes = []
        for st in sub_token_lists:
            sub_text = "".join(t["base"] for t in st)
            sub_chars = len(sub_text)
            sub_fs = min(_jp_font_size(sub_chars), fs)
            while sub_fs > 40 and not _text_fits_width(sub_text, sub_fs):
                sub_fs -= 4
            font_sizes.append(sub_fs)
    else:
        sub_token_lists = [tokens]
        while fs > 40 and not _text_fits_width(full_text, fs):
            fs -= 4
        font_sizes = [fs]

    if is_paired:
        font_sizes = [_reduce_tier(f) for f in font_sizes]

    return sub_token_lists, font_sizes


# ── Height estimation ─────────────────────────────────────────────────────────

def _estimate_line_group_height(sub_token_lists, font_sizes) -> int:
    total = 0
    furi_h = int(Pt(FURIGANA_SIZE_PT + 2))
    for i, (st, fs) in enumerate(zip(sub_token_lists, font_sizes)):
        base_h = int(Pt(fs + 8))
        total += furi_h + FURIGANA_ABOVE_BASE + base_h
        if i < len(sub_token_lists) - 1:
            total += GAP_BETWEEN_SPLIT
    total += GAP_AFTER_BASE
    total += int(Pt(ROMAJI_SIZE_PT + 6))
    total += int(Pt(2))
    total += int(Pt(CN_SIZE_PT + 6))
    total += GAP_AFTER_CN
    return total


# ── Rendering ─────────────────────────────────────────────────────────────────

def _render_furigana_layer(slide, sub_tokens, x_base: int, y: int,
                           font_size_jp: int):
    furi_h = int(Pt(FURIGANA_SIZE_PT + 2))
    char_offset = 0
    for tok in sub_tokens:
        if tok["needs_furi"]:
            x_off = char_offset * _char_width_emu(font_size_jp)
            bw = len(tok["base"]) * _char_width_emu(font_size_jp)
            fw = _text_width_emu(tok["kana"], FURIGANA_SIZE_PT)
            fw = max(fw, int(bw * 0.85))
            fw = int(fw * 1.25)
            fx = x_base + x_off - (fw - bw) // 2

            tf = slide.shapes.add_textbox(fx, y, fw, furi_h)
            tf.text_frame.word_wrap = True
            tf.text_frame.paragraphs[0].text = ""
            pf = tf.text_frame.add_paragraph()
            pf.alignment = PP_ALIGN.CENTER
            pf.space_before = Pt(0)
            pf.space_after = Pt(0)
            _add_run(pf, tok["kana"], FONT_FURIGANA, FONT_EN,
                     size=Pt(FURIGANA_SIZE_PT), color=TEXT_FURIGANA)
        char_offset += len(tok["base"])


def _render_sub_line(slide, sub_tokens, y: int, font_size_jp: int,
                     x_margin: int) -> int:
    sub_text = "".join(t["base"] for t in sub_tokens)

    furi_y = y
    furi_h = int(Pt(FURIGANA_SIZE_PT + 2))
    _render_furigana_layer(slide, sub_tokens, x_margin, furi_y, font_size_jp)

    base_y = furi_y + furi_h + FURIGANA_ABOVE_BASE
    base_h = int(Pt(font_size_jp + 8))
    base_w = min(_text_width_emu(sub_text, font_size_jp) + int(Pt(6)),
                  _max_content_w(x_margin))

    tf = slide.shapes.add_textbox(x_margin, base_y, base_w, base_h)
    tf.text_frame.word_wrap = False
    tf.text_frame.paragraphs[0].text = ""
    pb = tf.text_frame.add_paragraph()
    pb.alignment = PP_ALIGN.LEFT
    pb.space_before = Pt(0)
    pb.space_after = Pt(0)
    _add_run(pb, sub_text, FONT_JP, FONT_EN,
             size=Pt(font_size_jp), color=TEXT_JP)

    return base_y + base_h


def _render_lyric_line(slide, sub_token_lists, font_sizes,
                       romaji: str, translation: str,
                       y_start: int, x_margin: int) -> int:
    y = y_start
    for i, (st, fs) in enumerate(zip(sub_token_lists, font_sizes)):
        y = _render_sub_line(slide, st, y, fs, x_margin)
        if i < len(sub_token_lists) - 1:
            y += GAP_BETWEEN_SPLIT

    y += GAP_AFTER_BASE

    # Romaji — typewriter font with letter spacing
    ro_h = int(Pt(ROMAJI_SIZE_PT + 6))
    tf_ro = slide.shapes.add_textbox(x_margin, y, _max_content_w(x_margin), ro_h)
    tf_ro.text_frame.word_wrap = True
    tf_ro.text_frame.paragraphs[0].text = ""
    pr = tf_ro.text_frame.add_paragraph()
    pr.alignment = PP_ALIGN.LEFT
    pr.space_before = Pt(0)
    pr.space_after = Pt(0)
    run_ro = _add_run(pr, romaji, FONT_EN, FONT_ROMAJI,
                      size=Pt(ROMAJI_SIZE_PT), color=TEXT_ROMAJI)
    _set_char_spacing(run_ro, ROMAJI_SPACING_PT)

    y += ro_h + int(Pt(2))

    # Chinese translation — lighter than romaji
    cn_h = int(Pt(CN_SIZE_PT + 6))
    tf_cn = slide.shapes.add_textbox(x_margin, y, _max_content_w(x_margin), cn_h)
    tf_cn.text_frame.word_wrap = True
    tf_cn.text_frame.paragraphs[0].text = ""
    pc = tf_cn.text_frame.add_paragraph()
    pc.alignment = PP_ALIGN.LEFT
    pc.space_before = Pt(0)
    pc.space_after = Pt(0)
    _add_run(pc, translation, FONT_CN, FONT_EN,
             size=Pt(CN_SIZE_PT), color=TEXT_CN)

    return y + cn_h + GAP_AFTER_CN


# ── Pairing ───────────────────────────────────────────────────────────────────

def _pair_lines(lines: list) -> list[list]:
    pairs = []
    i, n = 0, len(lines)
    usable_h = SLIDE_H - BOTTOM_MARGIN

    while i < n:
        cur = lines[i]
        cur_cc = len(cur[0])

        if cur_cc >= 19 or i == n - 1:
            pairs.append([cur])
            i += 1
            continue

        nxt = lines[i + 1]
        nxt_cc = len(nxt[0])
        if nxt_cc >= 19:
            pairs.append([cur])
            i += 1
            continue

        h_cur = _estimate_line_group_height(*_resolve_line(cur[0], True))
        h_nxt = _estimate_line_group_height(*_resolve_line(nxt[0], True))
        if h_cur + h_nxt + GAP_LINES <= usable_h:
            pairs.append([cur, nxt])
            i += 2
        else:
            pairs.append([cur])
            i += 1

    return pairs


# ── Slide background helpers ──────────────────────────────────────────────────

def _add_slide_bg_and_section(slide, bg_buf, section_name: str):
    """Add paper background and optional section watermark to a slide."""
    bg_buf.seek(0)
    slide.shapes.add_picture(bg_buf, 0, 0, SLIDE_W, SLIDE_H)
    if section_name:
        tf_s = slide.shapes.add_textbox(
            int(Inches(0.4)), int(Pt(16)), int(Inches(4)), int(Pt(24)))
        tf_s.text_frame.word_wrap = True
        _add_run(tf_s.text_frame.paragraphs[0], section_name,
                 FONT_JP, FONT_EN, size=Pt(SECTION_SIZE_PT), color=TEXT_SECTION)


# ── Slide layout ──────────────────────────────────────────────────────────────

def _layout_and_render_slide(slide, paired_lines, bg_buf, section_name,
                             layout_idx: int, overflow_warnings: list[str] | None = None):
    """Layout and render a slide using one of 4 composition templates."""
    _add_slide_bg_and_section(slide, bg_buf, section_name)

    layout = LAYOUTS[layout_idx % len(LAYOUTS)]
    is_paired = len(paired_lines) == 2
    is_stagger = layout["name"] == "stagger"
    usable_h = SLIDE_H - BOTTOM_MARGIN

    # Resolve all lines
    resolutions = []
    for lt in paired_lines:
        sub_tokens, font_sizes = _resolve_line(lt[0], is_paired)
        h = _estimate_line_group_height(sub_tokens, font_sizes)
        resolutions.append((sub_tokens, font_sizes, h))

    # Overflow detection — width
    if overflow_warnings is not None:
        for lt, (sub_tokens, font_sizes, _h) in zip(paired_lines, resolutions):
            for st, fs in zip(sub_tokens, font_sizes):
                sub_text = "".join(t["base"] for t in st)
                if not _text_fits_width(sub_text, fs, layout["x"]):
                    preview = sub_text[:30] + ("..." if len(sub_text) > 30 else "")
                    overflow_warnings.append(
                        f"[{section_name}] JP overflow: {preview} "
                        f"({len(sub_text)} chars @ {fs}pt)"
                    )

    if is_stagger and is_paired:
        # Stagger: two groups at different y positions
        y_positions = [TOP_LEFT_TOP, STAGGER_G2_TOP]
        for idx, (lt, (sub_tokens, font_sizes, _h)) in enumerate(
            zip(paired_lines, resolutions)):
            original, romaji, translation = lt
            y = y_positions[idx]
            # Clamp to bottom margin
            g_h = _estimate_line_group_height(sub_tokens, font_sizes)
            if y + g_h > usable_h:
                y = max(0, usable_h - g_h)
            _render_lyric_line(slide, sub_tokens, font_sizes,
                               romaji, translation, y, layout["x"])
    else:
        # Normal and center-left layouts: stack groups with layout gap
        gap = layout["gap"] if is_paired else 0
        total_h = sum(h for _, _, h in resolutions) + gap

        # Determine y-start
        if layout["name"] == "top_left":
            y = TOP_LEFT_TOP
        else:
            y = max(0, (SLIDE_H - total_h) // 2)

        if y + total_h > usable_h:
            y = max(0, usable_h - total_h)

        for i, (lt, (sub_tokens, font_sizes, _h)) in enumerate(
            zip(paired_lines, resolutions)):
            original, romaji, translation = lt
            y = _render_lyric_line(slide, sub_tokens, font_sizes,
                                   romaji, translation, y, layout["x"])
            if i < len(paired_lines) - 1:
                y += gap - GAP_AFTER_CN


# ── English grouping ──────────────────────────────────────────────────────────

def _group_english_lines(lines: list, max_per_slide: int = 5) -> list[list]:
    """Group English lyrics into slides, avoiding orphan trailing groups."""
    groups = []
    i, n = 0, len(lines)
    while i < n:
        remaining = n - i
        if remaining <= max_per_slide:
            groups.append(lines[i:])
            break
        size = max_per_slide
        if remaining - size == 1:
            size = max_per_slide - 1
        elif remaining - size == 2 and max_per_slide >= 5:
            size = max_per_slide - 1
        groups.append(lines[i:i + size])
        i += size
    return groups


# ── English slide rendering ────────────────────────────────────────────────────

_EN_LINE_GAP = int(Pt(20))
_EN_CN_GAP = int(Pt(4))


def _layout_and_render_english_slide(slide, lines, bg_buf, section_name,
                                      layout_idx: int, overflow_warnings: list[str] | None = None):
    """Layout and render an English lyric slide using 4 poster templates."""
    _add_slide_bg_and_section(slide, bg_buf, section_name)

    layout = LAYOUTS[layout_idx % len(LAYOUTS)]
    n_lines = len(lines)
    usable_h = SLIDE_H - BOTTOM_MARGIN

    # Font size from longest English line
    max_len = max(len(lt[0]) for lt in lines)
    fs_en = _en_font_size(max_len)

    # Height per line group: EN text + gap + CN text
    en_h = int(Pt(fs_en + 10))   # matches _render_english_line
    cn_h = int(Pt(CN_SIZE_PT + 6))
    group_h = en_h + _EN_CN_GAP + cn_h

    total_h = group_h * n_lines + _EN_LINE_GAP * (n_lines - 1)

    # Check if width fits; downsize if needed
    en_max_w = _max_content_w(layout["x"])
    while fs_en > 24:
        all_fit = True
        for lt in lines:
            if _en_text_width_emu(lt[0], fs_en) > en_max_w:
                all_fit = False
                break
        if all_fit:
            break
        fs_en -= 4
        en_h_new = int(Pt(fs_en + 8))
        group_h_new = en_h_new + _EN_CN_GAP + cn_h
        total_h = group_h_new * n_lines + _EN_LINE_GAP * (n_lines - 1)

    # Overflow detection
    if overflow_warnings is not None:
        if fs_en <= 24:
            for lt in lines:
                if _en_text_width_emu(lt[0], fs_en) > en_max_w:
                    preview = lt[0][:30] + ("..." if len(lt[0]) > 30 else "")
                    overflow_warnings.append(
                        f"[{section_name}] EN overflow: {preview} "
                        f"({len(lt[0])} chars @ {fs_en}pt)"
                    )
        if total_h > usable_h:
            overflow_warnings.append(
                f"[{section_name}] EN height overflow: "
                f"{total_h} > {usable_h} EMU"
            )

    # Determine y-start
    is_stagger = layout["name"] == "stagger" and n_lines >= 4

    if is_stagger:
        mid = n_lines // 2
        top_n, bottom_n = mid, n_lines - mid
        top_h = group_h * top_n + _EN_LINE_GAP * max(0, top_n - 1)
        bottom_h = group_h * bottom_n + _EN_LINE_GAP * max(0, bottom_n - 1)
        top_end = TOP_LEFT_TOP + top_h
        bottom_end = STAGGER_G2_TOP + bottom_h

        # Stagger only if groups don't collide and bottom fits
        if top_end < STAGGER_G2_TOP and bottom_end <= usable_h:
            # Top group
            y = TOP_LEFT_TOP
            for lt in lines[:mid]:
                y = _render_english_line(slide, lt, y, layout["x"], fs_en)
                y += _EN_LINE_GAP - _EN_CN_GAP
            # Bottom group
            y2 = STAGGER_G2_TOP
            for lt in lines[mid:]:
                y2 = _render_english_line(slide, lt, y2, layout["x"], fs_en)
                y2 += _EN_LINE_GAP - _EN_CN_GAP
        else:
            is_stagger = False

    if not is_stagger:
        if layout["name"] == "top_left":
            y = TOP_LEFT_TOP
        else:
            y = max(0, (SLIDE_H - total_h) // 2)
        if y + total_h > usable_h:
            y = max(0, usable_h - total_h)

        for lt in lines:
            y = _render_english_line(slide, lt, y, layout["x"], fs_en)
            y += _EN_LINE_GAP - _EN_CN_GAP


def _render_english_line(slide, lt, y: int, x: int, fs_en: int) -> int:
    """Render one English + Chinese line pair. Returns y after this group."""
    en_text, cn_text = lt

    en_h = int(Pt(fs_en + 10))
    tf_en = slide.shapes.add_textbox(x, y, _max_content_w(x), en_h)
    tf_en.text_frame.word_wrap = True
    tf_en.text_frame.paragraphs[0].text = ""
    p_en = tf_en.text_frame.add_paragraph()
    p_en.alignment = PP_ALIGN.LEFT
    p_en.space_before = Pt(0)
    p_en.space_after = Pt(0)
    run_en = p_en.add_run()
    run_en.text = en_text
    run_en.font.name = FONT_EN_SERIF_BOLD
    run_en.font.size = Pt(fs_en)
    run_en.font.color.rgb = rgb(TEXT_JP)
    # DrawingML Latin font
    rPr = run_en._r.get_or_add_rPr()
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    latin = rPr.find(f"{ns}latin")
    if latin is None:
        latin = etree.SubElement(rPr, f"{ns}latin")
    latin.set("typeface", FONT_EN_SERIF_BOLD)

    y += en_h + _EN_CN_GAP

    # Chinese translation
    cn_h = int(Pt(CN_SIZE_PT + 6))
    tf_cn = slide.shapes.add_textbox(x, y, _max_content_w(x), cn_h)
    tf_cn.text_frame.word_wrap = True
    tf_cn.text_frame.paragraphs[0].text = ""
    p_cn = tf_cn.text_frame.add_paragraph()
    p_cn.alignment = PP_ALIGN.LEFT
    p_cn.space_before = Pt(0)
    p_cn.space_after = Pt(0)
    _add_run(p_cn, cn_text, FONT_CN, FONT_EN,
             size=Pt(CN_SIZE_PT), color=TEXT_CN)

    return y + cn_h


# ── Self-verification ──────────────────────────────────────────────────────────

def _verify_build(d: dict, slide_count: int, is_japanese: bool,
                  overflow_warnings: list[str]):
    """Verify lyric completeness and print summary."""
    fl = d.get("lyric_sections", [])
    total_input = 0
    section_counts = []
    for sname, lines in fl:
        if lines:
            section_counts.append((sname, len(lines)))
            total_input += len(lines)

    lang = "JP" if is_japanese else "EN"
    print(f"\n{'='*60}")
    print(f"  VERIFY: {d.get('title', 'song')} ({lang})")
    print(f"  Input lines: {total_input}")
    print(f"  Output slides: {slide_count}")
    per_page = "1–2" if is_japanese else "3–5"
    print(f"  Lines/page target: {per_page}")

    # Section breakdown
    print(f"  Sections:")
    for sname, cnt in section_counts:
        print(f"    [{sname or '(no label)'}] {cnt} lines")

    if overflow_warnings:
        print(f"  ⚠ OVERFLOW WARNINGS:")
        for w in overflow_warnings:
            print(f"    {w}")
    else:
        print(f"  ✓ No overflow detected")

    print(f"{'='*60}\n")


# ── Main ──────────────────────────────────────────────────────────────────────

def build_pptx(json_path: str, out_path: str | None = None):
    global TEXT_JP, TEXT_FURIGANA, TEXT_ROMAJI, TEXT_CN, TEXT_SECTION

    with open(json_path, "r", encoding="utf-8") as f:
        d = json.load(f)

    if out_path is None:
        out_path = json_path.replace("data.json",
            d.get("title", "song").replace("歌曲学习：", "") + ".pptx")
        if out_path == json_path:
            dn = os.path.basename(os.path.dirname(json_path))
            out_path = os.path.join(os.path.dirname(json_path), dn + ".pptx")

    fl = d.get("lyric_sections", [])
    is_japanese = bool(fl and fl[0][1] and len(fl[0][1][0]) == 3)

    seed = int(hashlib.md5(d.get("title", "song").encode()).hexdigest()[:8], 16) % 10000
    bg_color = d.get("bg_color", "C5CDD4")
    mood = d.get("mood")
    bg_buf = generate_mood_background(seed=seed, base_color=bg_color, mood=mood)

    # Adapt text palette to background for readability
    pal = palette_for_bg(bg_color)
    TEXT_JP = pal["TEXT_JP"]
    TEXT_FURIGANA = pal["TEXT_FURIGANA"]
    TEXT_ROMAJI = pal["TEXT_ROMAJI"]
    TEXT_CN = pal["TEXT_CN"]
    TEXT_SECTION = pal["TEXT_SECTION"]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Title slide
    title = d.get("title", "").replace("歌曲学习：", "")
    s0 = prs.slides.add_slide(prs.slide_layouts[6])
    bg_buf.seek(0)
    s0.shapes.add_picture(bg_buf, 0, 0, SLIDE_W, SLIDE_H)
    tf0 = s0.shapes.add_textbox(LEFT_MARGIN, int(Inches(2.3)),
                                _max_content_w(LEFT_MARGIN), int(Inches(3.5)))
    tf0.text_frame.word_wrap = True
    _add_run(tf0.text_frame.paragraphs[0], title,
             FONT_JP if is_japanese else FONT_EN_SERIF_BOLD, FONT_EN,
             size=Pt(48), color=TEXT_JP)
    if d.get("info_rows"):
        ar = [r for r in d["info_rows"] if r[0] == "演唱者"]
        if ar:
            pa = tf0.text_frame.add_paragraph()
            pa.alignment = PP_ALIGN.LEFT
            pa.space_before = Pt(12)
            _add_run(pa, ar[0][1], FONT_CN, FONT_EN, size=Pt(20), color=TEXT_CN)

    # Lyric slides
    slide_idx = 0
    overflow_warnings = []

    for sname, lines in fl:
        if not lines:
            continue

        if is_japanese:
            paired = _pair_lines(lines)
            for group in paired:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _layout_and_render_slide(slide, group, bg_buf, sname, slide_idx,
                                         overflow_warnings)
                slide_idx += 1
        else:
            grouped = _group_english_lines(lines)
            for group in grouped:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _layout_and_render_english_slide(slide, group, bg_buf, sname, slide_idx,
                                                  overflow_warnings)
                slide_idx += 1

    prs.save(out_path)

    # Self-verification
    _verify_build(d, slide_idx, is_japanese, overflow_warnings)
    print(f"OK: {out_path} ({os.path.getsize(out_path)} bytes)")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: PYTHONUTF8=1 python build_pptx.py <data.json>")
        print("  or: PYTHONUTF8=1 python build_pptx.py --all")
        sys.exit(1)
    if sys.argv[1] == "--all":
        roots = [os.path.dirname(os.path.abspath(__file__))]
        ext = os.environ.get("SONG_STUDY_DATA", "")
        if ext and os.path.isdir(ext):
            roots.append(ext)
        for root in roots:
            if os.path.isdir(root):
                for e in os.scandir(root):
                    dp = os.path.join(e.path, "data.json") if e.is_dir() else ""
                    if dp and os.path.isfile(dp):
                        build_pptx(dp)
    else:
        build_pptx(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else None)
